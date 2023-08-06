#!/usr/bin/python
# -*- coding: UTF-8 -*-
#载入必要的模块
import pygame
from datetime import datetime
from faker import Factory
from faker import Faker
from docx import Document
from docx.shared import Inches
from multiprocessing import Pool
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
from pptx.enum.chart import XL_LEGEND_POSITION
from reportlab.pdfgen import canvas
from reportlab.platypus.tables import Table, TableStyle
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph,Frame
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import Image as platImage
from PIL import Image
from reportlab.lib import colors
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
import time
import random
import socket
import getopt
import zipfile
import sys
import os
import xlwt
import xlrd
import string
import struct
import _winreg
import shutil
import subprocess
from PIL import ImageGrab
import win32api
import win32con


reload(sys)
sys.setdefaultencoding('utf8')

class office:

    def hello(self):
        print('hello')

    def timeget(self):
        timeStamp = time.time()
        timeArray = time.localtime(timeStamp)
        otherStyleTime = time.strftime("%Y%m%d%H%M%S", timeArray) 
        return otherStyleTime

    def mkdir(self,path):
        path=path.strip()
        path=path.rstrip("\\")
        isExists=os.path.exists(path)
        if not isExists:
            os.makedirs(path.decode('utf-8'))
            print path + u' 创建成功'
            return path
        else:
            print path + u' 目录已存在'
            return path

    def toText(self,f):
        faker = Faker("zh_CN")
        f.writelines("王宝强")
        f.writelines("\n")
        f.writelines("苏绍辉")
        f.writelines("\n")
        f.writelines("姓名: ")
        f.writelines(faker.name())
        f.writelines("\n")
        f.writelines("电话号码: ")
        f.writelines(faker.phone_number())
        f.writelines("\n")
        f.writelines("18位身份证号: ")
        f.writelines(faker.ssn())
        f.writelines("\n")
        f.writelines("信用卡卡号: ")
        f.writelines(faker.credit_card_number(card_type=None))
        f.writelines("\n")
        f.writelines("地址: ")
        f.writelines(faker.address())

    def genfaker(self,num):
        faker = Faker("zh_CN")
        file = 'Result\\card.txt'
        if os.path.exists(file):
            os.remove(file)
            pass
        f = open(file,'a')
        i = 0
        while(i < num):
            f.writelines("姓名: ")
            f.writelines(faker.name())
            f.writelines("\n")
            f.writelines("电话号码: ")
            f.writelines(faker.phone_number())
            f.writelines("\n")
            f.writelines("18位身份证号: ")
            f.writelines(faker.ssn())
            f.writelines("\n")
            f.writelines("信用卡卡号: ")
            f.writelines(faker.credit_card_number(card_type=None))
            f.writelines("\n")
            f.writelines("地址: ")
            f.writelines(faker.address())
            f.writelines("\n----------------------------------------\n")
            #f.writelines(["phone: " + faker.phone_number() + "\n","user sn: " + faker.ssn() + "\n","bank card: " + faker.credit_card_number(card_type=None) + "\n"])
            i = i + 1
        f.close()
        return True

    def genContent(self,num):
        o = office()
        fix = o.timeget()
        spath = "Result\\txt\\%s" % fix
        path = o.mkdir(spath)
        size = random.randint(1,100)

        faker = Faker("zh_CN")
        file = '%s\\content_%s.txt' % (path,num)
        s = faker.text(max_nb_chars=2048)
        if os.path.exists(file):
            os.remove(file)
            pass
        f = open(file,'a')
        i = 0
        #while(i < 185*size):#循环185次刚好是1M大小文件,循环2次刚好是12K大小
        while(i < 20*size):#循环185次刚好是1M大小文件,循环2次刚好是12K大小
            f.write(s)
            i = i + 1
        o.toText(f)
        f.close()
        return True

    def CreateContent(self,n):
        m = 0
        o = office()
        while m < n:
            o.genContent(m)
            m = m + 1
        return True
    
    def CreateExcel(self,text,s):
        try:
            i = 0
            z = 3
            style = xlwt.XFStyle()
            font = xlwt.Font()
            font.name = 'Times New Roman'
            font.bold = True
            style.font = font
            filename = xlwt.Workbook(encoding = 'UTF-8')
            while i < s:
                t = 0
                x = random.randint(0, 10)
                y = random.randint(0, 10)
                sheet = filename.add_sheet("test %s" % i)
                #while t < 100:
                while t < 10001:
                    sheet.write(x,y,text,style)
                    x = x + 1
                    t = t + 1
                i = i + 1
            sheet = filename.add_sheet("test %s" % s)
            while z < 10:
                salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
                sheet.write(2,z - 3,"王宝强-苏绍辉" + salt,style)
                z = z + 1
            filename.save("Result\\dlp.xls")
            return True
        except Exception,e:
            print(str(e))
            return False
    def CreateExcelNum(self,text,s,num,r,c):
        path = "Result\\excel"
        o.mkdir(path)
        try:
            i = 0
            z = 3
            style = xlwt.XFStyle()
            font = xlwt.Font()
            font.name = 'Times New Roman'
            font.bold = True
            style.font = font
            filename = xlwt.Workbook(encoding = 'UTF-8')
            while i < s:
                t = 0
                wa =0
                qc = 0
                x = random.randint(0, 10)
                #y = random.randint(0, 10)
                sheet = filename.add_sheet("test %s" % i)
                #while t < 100:
                while t < 100:
                    while wa < r:                       # r 代表行
                        for qc in xrange(0,c):            #c 代表列
                            sheet.write(wa,qc,text,style)
                            qc = qc + 1
                        wa = wa + 1
                    t = t + 1
                i = i + 1
            sheet = filename.add_sheet("test %s" % s)
            while z < 10:
                salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
                sheet.write(2,z - 3,"王宝强-苏绍辉" + salt,style)
                z = z + 1
            filename.save("Result\\excel\\dlp%s.xls" % num)
            return True
        except Exception,e:
            print(str(e))
            return False
    def CreateExcelE(self,text,s,n,r,c):
        m = 0
        o = office()
        while m < n:
            o.CreateExcelNum(text,s,m,r,c)
            m = m + 1
        return True

    def CreateExcelE1(self,n):
        o = office()
        m = 0
        faker = Faker("zh_CN")
        text = faker.text(max_nb_chars=255)

        while m < n:
            s = random.randint(50,100)
            o.CreateExcelNum(text,s,m,100,50)
            m = m + 1
        return True

    def CreateWord(self,s):
        i = 0
        z = 0
        document = Document()

        while i < s:
            document.add_heading(u'文档标题', 0)
            p = document.add_paragraph(u'这是一个自然段 ')
            p.add_run('bold').bold = True
            p.add_run(u' 还有 ')
            p.add_run('italic.').italic = True

            document.add_heading(u'1级别标题', level=1)
            document.add_paragraph(u'引用', style='IntenseQuote')

            document.add_paragraph(u'符号列表', style='ListBullet')
            document.add_paragraph(u'数字列表t', style='ListNumber')
            document.add_paragraph(u'我的微信:')
            document.add_picture('data\\tupian.jpg', width=Inches(3.25))

            table = document.add_table(rows=3, cols=3)
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = u'第一列'
            hdr_cells[1].text = u'第二列'
            hdr_cells[2].text = u'第三列'

            hdr_cells = table.rows[1].cells
            hdr_cells[0].text = '1'
            hdr_cells[1].text = '21'
            hdr_cells[2].text = 'qwertyuiop'

            hdr_cells = table.rows[2].cells
            hdr_cells[0].text = '2'
            hdr_cells[1].text = '43'
            hdr_cells[2].text = 'asdfghjkl'
            i = i + 1
       
        document.add_heading(u'我是关键字', 0)
        while z < 10:
            faker = Faker("zh_CN")
            text = faker.text(max_nb_chars=2048)
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            document.add_paragraph(u'王宝强-苏绍辉' + text + salt)
            z = z + 1

        document.add_page_break()

        document.save("Result\\dlp.docx")

        return True
    def CreateWordNum(self,num):
        path = "Result\\word"
        o.mkdir(path)
        i = 0
        z = 0
        document = Document()
        while i < 100:
            document.add_heading(u'文档标题', 0)
            p = document.add_paragraph(u'这是一个自然段 ')
            p.add_run('bold').bold = True
            p.add_run(u' 还有 ')
            p.add_run('italic.').italic = True

            document.add_heading(u'1级别标题', level=1)
            document.add_paragraph(u'引用', style='IntenseQuote')

            document.add_paragraph(u'符号列表', style='ListBullet')
            document.add_paragraph(u'数字列表t', style='ListNumber')
            document.add_paragraph(u'我的微信:')
            document.add_picture('data\\tupian.jpg', width=Inches(3.25))

            table = document.add_table(rows=3, cols=3)
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = u'第一列'
            hdr_cells[1].text = u'第二列'
            hdr_cells[2].text = u'第三列'

            hdr_cells = table.rows[1].cells
            hdr_cells[0].text = '1'
            hdr_cells[1].text = '21'
            hdr_cells[2].text = 'qwertyuiop'

            hdr_cells = table.rows[2].cells
            hdr_cells[0].text = '2'
            hdr_cells[1].text = '43'
            hdr_cells[2].text = 'asdfghjkl'
            i = i + 1
       
        document.add_heading(u'我是关键字', 0)
        s = random.randint(0,50)
        while z < s:
            faker = Faker("zh_CN")
            text = faker.text(max_nb_chars=20480)
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            document.add_paragraph(u'王宝强-苏绍辉' + text + salt)
            z = z + 1

        document.add_page_break()

        document.save("Result\\word\\dlp%s.docx" % num)

        return True
    def CreateWordE(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreateWordNum(m)
            m = m + 1
        return True
    def CreateWordE1(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreateWordNum(m)
            m = m + 1
        return True

    def CreatePpt(self):
        i = 0
        z = 0
        # 创建幻灯片 ------
        prs = Presentation()
        while i < 10:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # 定义图表数据 ---------------------
            chart_data = ChartData()
            chart_data.categories = ['East', 'West', 'Midwest']
            chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
         
            # 将图表添加到幻灯片 --------------------
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            i = i + 1

        #title_slide_layout1 = prs.slide_layouts[0]
        #slide1 = prs.slides.add_slide(title_slide_layout1)
         
        #title1 = slide1.shapes.title
        #subtitle1 = slide1.placeholders[1]
         
        # 设置标题和副标题
        #title1.text = "我是关键字"
        #salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #subtitle1.text = "王宝强-苏绍辉" + salt

        s = random.randint(0,50)
        while z < s:
            title_slide_layout1 = prs.slide_layouts[0]
            slide1 = prs.slides.add_slide(title_slide_layout1)
         
            title1 = slide1.shapes.title
            subtitle1 = slide1.placeholders[1]
            faker = Faker("zh_CN")
            text = faker.text(max_nb_chars=2048)
            subtitle1.text = text
            title1.text = "我是关键字"
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            subtitle1.text = "王宝强-苏绍辉" + salt
            z = z + 1
            
        title1.text = "我是关键字"
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        subtitle1.text = "王宝强-苏绍辉" + salt

        prs.save('Result\\dlp.pptx')

        return True

    def CreatePptNum(self,num):
        path = "Result\\ppt"
        o.mkdir(path)
        i = 0
        z = 0
        # 创建幻灯片 ------
        prs = Presentation()
        while i < 10:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # 定义图表数据 ---------------------
            chart_data = ChartData()
            chart_data.categories = ['East', 'West', 'Midwest']
            chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
         
            # 将图表添加到幻灯片 --------------------
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            i = i + 1

        #title_slide_layout1 = prs.slide_layouts[0]
        #slide1 = prs.slides.add_slide(title_slide_layout1)
         
        #title1 = slide1.shapes.title
        #subtitle1 = slide1.placeholders[1]
         
        # 设置标题和副标题
        #title1.text = "我是关键字"
        #salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #subtitle1.text = "王宝强-苏绍辉" + salt


        s = random.randint(0,50)
        while z < s:
            title_slide_layout1 = prs.slide_layouts[0]
            slide1 = prs.slides.add_slide(title_slide_layout1)
         
            title1 = slide1.shapes.title
            subtitle1 = slide1.placeholders[1]
            faker = Faker("zh_CN")
            text = faker.text(max_nb_chars=2048)
            subtitle1.text = text
            title1.text = "我是关键字"
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            subtitle1.text = "王宝强-苏绍辉" + salt
            z = z + 1

        prs.save('Result\\ppt\\dlp%s.pptx' % num)

        return True

    def CreatePptE(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreatePptNum(m)
            m = m + 1
        return True
    def CreatePptE1(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreatePptNum(m)
            m = m + 1
        return True

    def CreateTxt(self,s):
        first = []
        second = []
        f = open('Result\\dlp.txt','w')
        with open('file\\ftp.txt', 'r') as f1:
            for line in f1:
                line = line.strip()
                first.append(line)
            f1.close()
        with open('file\\ftp.txt', 'r') as f2:
            for line2 in f2:
                line2 = line2.strip()
                second.append(line2)
            f2.close()
        for i in range(0,s):
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            result = first[i] + '\t' + second[i] + '\n'
            f.write(result)
        f.write(salt)
        f.close()
        return True

    def CreateTxtE(self,s,n):
        first = []
        second = []
        x = 0
        while x < n:
            f = open('Result\\dlp%s.txt' % x,'w')
            with open('file\\ftp.txt', 'r') as f1:
                for line in f1:
                    line = line.strip()
                    first.append(line)
                f1.close()
            with open('file\\ftp.txt', 'r') as f2:
                for line2 in f2:
                    line2 = line2.strip()
                    second.append(line2)
                f2.close()
            for i in range(0,s):
                salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
                result = first[i] + '\t' + second[i] + '\n'
                f.write(result)
            f.write(salt)
            f.close()
            x = x + 1
        return True

    def CreateImage(self):
        pygame.init()
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #待转换文字
        text = u"woshi啊实打实大家卡萨丁比较卡失败的健康萨比\
        大数据科技部科技表示大空间的巴萨健康大鼠标接口\
        12372193721983198diasdhiuadslkandklasnbdklasd王宝强" + salt
        #设置字体和字号
        font = pygame.font.SysFont('Microsoft YaHei', 64)
        #渲染图片，设置背景颜色和字体样式,前面的颜色是字体颜色
        ftext = font.render(text, True, (65, 83, 130),(255, 255, 255))
        #保存图片
        pygame.image.save(ftext,"Result\\dlp.png")#图片保存地址
        return True
    def CreateImageNum(self,num):
        pygame.init()
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #待转换文字
        text = u"woshi啊实打实大家卡萨丁比较卡失败的健康萨比\
        大数据科技部科技表示大空间的巴萨健康大鼠标接口\
        12372193721983198diasdhiuadslkandklasnbdklasd王宝强" + salt
        #设置字体和字号
        font = pygame.font.SysFont('Microsoft YaHei', 64)
        #渲染图片，设置背景颜色和字体样式,前面的颜色是字体颜色
        ftext = font.render(text, True, (65, 83, 130),(255, 255, 255))
        #保存图片
        pygame.image.save(ftext,"Result\\dlp%s.png" % num)#图片保存地址
        return True
    def CreateImageE(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreateImageNum(m)
            m = m + 1
        return True

    def CreatePdf(self,s):
        #faker = Faker("zh_CN")
        x = 0
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #pdfmetrics.registerFont(TTFont('hei', 'hei.TTF'))
        #import testSubFun
        #testSubFun.testSubFunc('first')
        #设置页面大小
        c = canvas.Canvas('Result\\dlp.pdf',pagesize=A4)
        xlength,ylength = A4
        #print('width:%d high:%d'%(xlength,ylength))
        #c.line(1,1,ylength/2,ylength)
        #设置文字类型及字号
        #c.setFont('hei',20)
        #生成一个table表格
        atable = [[1,2,3,4,5,6,7,8],[11,12,13,14,15,16,17,18]]
        t = Table(atable,50,20)
        t.setStyle(TableStyle([('ALIGN',(0,0),(3,1),'CENTER'),
                               ('INNERGRID',(0,0),(-1,-1),0.25,colors.black),
                               ('BOX',(0,0),(-1,-1),0.25,colors.black)]))
        textOb = c.beginText(1,ylength-10)
        indexVlaue = 0
        while(indexVlaue < ylength):
            #textStr = '''wo shi tupian---wo shi tupian--wo shi tupian--wo shi tupian%d'''%indexVlaue + salt
            textStr = '''wo shi tupian---wo shi tupian--wo shi tupian--wo shi tupian--苏绍辉%d'''%indexVlaue + salt
            #print('nextline,nextline%d'%indexVlaue)
            textOb.textLine(textStr)
            indexVlaue = indexVlaue + 1
            break
        c.drawText(textOb)
        
        #简单的图片载入
        imageValue = 'data\\dlp.png'
        c.drawImage(imageValue,97,97,650,650)
        #c.drawImage('file\\dlp.png',50,50,50,50)
        t.split(0,0)
        t.drawOn(c,100,1)
        c.showPage()
        #换页的方式不同的showPage
        while x < s:
            imageValue = 'data\\dlp.png'
            c.drawImage(imageValue,97,97,650,650)
            c.drawString(0,0,'tupian%s' % x)
            c.showPage()
            x = x + 1
        c.save()

    def CreatePdfNum(self,s,num):
        x = 0
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #pdfmetrics.registerFont(TTFont('hei', 'hei.TTF'))
        #import testSubFun
        #testSubFun.testSubFunc('first')
        #设置页面大小
        path = "Result\\pdf"
        o.mkdir(path)
        c = canvas.Canvas('Result\\pdf\\dlp%s.pdf' % num,pagesize=A4)
        xlength,ylength = A4
        #print('width:%d high:%d'%(xlength,ylength))
        #c.line(1,1,ylength/2,ylength)
        #设置文字类型及字号
        #c.setFont('hei',20)
        #生成一个table表格
        atable = [[1,2,3,4,5,6,7,8],[11,12,13,14,15,16,17,18]]
        t = Table(atable,50,20)
        t.setStyle(TableStyle([('ALIGN',(0,0),(3,1),'CENTER'),
                               ('INNERGRID',(0,0),(-1,-1),0.25,colors.black),
                               ('BOX',(0,0),(-1,-1),0.25,colors.black)]))
        textOb = c.beginText(1,ylength-10)
        indexVlaue = 0
        while(indexVlaue < ylength):
            #textStr = u'''wo shi tupian---wo shi tupian--wo shi tupian--wo shi tupian%d'''%indexVlaue + salt
            textStr = '''wo shi tupian---wo shi tupian--wo shi tupian--wo shi tupian--苏绍辉%d'''%indexVlaue + salt
            #print('nextline,nextline%d'%indexVlaue)
            textOb.textLine(textStr)
            indexVlaue = indexVlaue + 1
            break
        c.drawText(textOb)
        
        #简单的图片载入
        imageValue = 'data\\dlp.png'
        c.drawImage(imageValue,97,97,650,650)
        #c.drawImage('file\\dlp.png',50,50,50,50)
        t.split(0,0)
        t.drawOn(c,100,1)
        c.showPage()
        #换页的方式不同的showPage
        while x < s:
            imageValue = 'data\\dlp.png'
            c.drawImage(imageValue,97,97,650,650)
            c.drawString(0,0,'tupian%s' % x)
            c.showPage()
            x = x + 1
        c.save()
    def CreatePdfE(self,n):
        m = 0
        o = office()
        while m < n:
            s = random.randint(500,1000)
            o.CreatePdfNum(s,m)
            m = m + 1
        return True

    def GenSizeFile(self,fileSize):
        #file path
        filePath="Result\\dlpsize.txt"
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        
        # 生成固定大小的文件
        # date size
        #ds = 0
        with open(filePath, "w") as f:
            '''
            while ds < fileSize:
                f.write(str(round(random.uniform(-1000, 1000),2)))
                f.write("\n")
                ds = os.path.getsize(filePath)
            print ds
            '''
            f.write('王宝强%s' % salt)
            f.seek(1024*1024*fileSize)
            f.write('\x00')
            f.write('\n')
            #f.write('王宝强')
            f.close()
        return True
        # print(os.path.getsize(filePath))

    # 支持文件类型
    # 用16进制字符串的目的是可以知道文件头是多少字节
    # 各种文件头的长度不一样，少半2字符，长则8字符
    def typeList(self):
        return {
            "574C4653": WLFS
            }
 
    # 字节码转16进制字符串
    def bytes2hex(self,bytes):
        num = len(bytes)
        hexstr = u""
        for i in range(num):
            t = u"%x" % bytes[i]
            if len(t) % 2:
                hexstr += u"0"
            hexstr += t
        return hexstr.upper()
    '''

    # 字节码转16进制字符串
    def bytes2hex(self,bytes):
        print(u'关键码转码……')
        num = len(bytes)
        hexstr = u""
        for i in range(num):
            t = u"%x" % bytes[i]
            if len(t) % 2:
                hexstr += u"0"
            hexstr += t
        return hexstr.upper()
    '''
    def get_desktop(self):
        key = _winreg.OpenKey(_winreg.HKEY_CURRENT_USER,r'Software\Microsoft\Windows\CurrentVersion\Explorer\Shell Folders')
        return _winreg.QueryValueEx(key, "Desktop")[0]
    '''
    # 获取文件类型
    def filetype(self,filename):
        app = office()
        print(u'读文件二进制码中……')
        binfile = open(filename, 'rb') # 必需二制字读取
        print(u'提取关键码……')
        bins = binfile.read(20) #提取20个字符
        binfile.close() #关闭文件流
        bins = app.bytes2hex(bins) #转码
        bins = bins.lower()#小写
        print(bins)
        tl = app.typeList()#文件类型
        ftype = 'unknown'
        print(u'关键码比对中……')
        for hcode in tl.keys():
            lens = len(hcode) # 需要的长度
            if bins[0:lens] == hcode:
                ftype = tl[hcode]
                break
        if ftype == 'unknown':#全码未找到，优化处理，码表取5位验证
            bins = bins[0:5]
            for hcode in tl.keys():
                if len(hcode) > 5 and bins == hcode[0:5]:
                    ftype = tl[hcode]
                    break
        return ftype
    ''' 
    # 获取文件类型
    def filetype(self,filename):
        app = office()
        binfile = open(filename, 'rb') # 必需二制字读取
        tl = app.typeList()
        ftype = 'unknown'
        for hcode in tl.keys():
            numOfBytes = len(hcode) / 2 # 需要读多少字节
            binfile.seek(0) # 每次读取都要回到文件头，不然会一直往后读取
            hbytes = struct.unpack_from("B"*numOfBytes, binfile.read(numOfBytes)) # 一个 "B"表示一个字节
            f_hcode = app.bytes2hex(hbytes)
            if f_hcode == hcode:
                ftype = tl[hcode]
                break
        binfile.close()
        return ftype

    def enzip(self):
        target_dir = 'Result\\'
        target = target_dir + time.strftime('%Y%m%d%H%M%S') + '.zip'
        f = zipfile.ZipFile(target,'w',zipfile.ZIP_DEFLATED)
        startdir = "Result"
        for dirpath, dirnames, filenames in os.walk(startdir):
            for filename in filenames:
                f.write(os.path.join(dirpath,filename))
        print 'Sucessful backup to', target
        f.close()

    def enrar(self):
        source = ['Result']
        target_dir = 'Result\\'
        target = target_dir + time.strftime('%Y%m%d%H%M%S') + '.rar'
        zip_command = "tool\\WinRAR\\Rar.exe a %s %s" % (target, ' '.join(source))
        if os.system(zip_command) == 0:
            print 'Sucessful backup to', target
        else:
            print 'Backup Failed'

    def logtodis(self,text):
        logpath = 'log\\rundiscovery.txt'
        f = open(logpath,'a+')
        f.write(text)
        f.write('---rundiscovery:\t%s\r\n' %datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
        f.close()

    def rundiscovery(self):
        o = office()
        path = "Result"
        isExists = os.path.exists(path)
        if not isExists:
            o.mkdir(path)
        else:
            shutil.rmtree(path)
            time.sleep(15)
            o.mkdir(path)
        
        print "START------>" + datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        o.CreatePdfE(40)
        time.sleep(5)
        o.logtodis("PDF")
        o.CreateExcelE1(9)
        time.sleep(5)
        o.logtodis("EXCEL")
        o.CreateContent(200)
        time.sleep(5)
        o.logtodis("TXT")
        o.CreateWordE1(9)
        time.sleep(5)
        o.logtodis("WORD")
        o.CreatePptE1(40)
        time.sleep(5)
        o.logtodis("PPT")
        o.enzip()
        time.sleep(5)
        o.logtodis("ZIP")
        o.enrar()
        o.logtodis("RAR")
        print "STOP------>" + datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    def test(self):
        i = 0
        path = os.getcwd()+ '\\'
        time.sleep(3)
        starttime = time.time()
        for y in range(18):
            try:
                o = office()
                o.rundiscovery()
                time.sleep(3)
                f = timeget()
                print 'rundiscovery' + '  '+ str(i)+ '  ' + f

                i = i + 1
            except Exception,ex:
                print ex
            if y == 17:
                break
            time.sleep(3600)
        print 'over'
        win32api.keybd_event(win32con.VK_LWIN,0,0,0)
        win32api.keybd_event(68,0,0,0)
        win32api.keybd_event(68,0,win32con.KEYEVENTF_KEYUP,0)
        win32api.keybd_event(win32con.VK_LWIN,0,win32con.KEYEVENTF_KEYUP,0)          
        subprocess.Popen('taskmgr')
        time.sleep(4)
        im = ImageGrab.grab()
        im.save('c:\\12.png')      
        time.sleep(120)
        try:
            os.system('taskkill /IM 360tray.exe /F')
            os.system('taskkill /IM 360EntClient.exe /F')
            os.system('taskkill /IM DLPMain.exe /F')
            os.system('taskkill /IM DlpCenter.exe /F')
            os.system('taskkill /IM dlpmain32.exe /F')
            os.system('taskkill /IM DriverControler.exe /F')
            os.system('taskkill /IM DriverControler.exe /F')
        except Exception,ex:
            print ex
    
#end--------------------------------------------------------

if __name__ == "__main__":
    o = office()
    o.test()