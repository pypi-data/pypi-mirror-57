#!/usr/bin/python
# -*- coding: UTF-8 -*-
#载入必要的模块

from datetime import datetime
from faker import Factory
from faker import Faker
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
from pptx.enum.chart import XL_LEGEND_POSITION
import time
import random
import sys
import os
import shutil
import ConfigParser
import string


reload(sys)
sys.setdefaultencoding('utf8')

class office:

    def hello(self):
        print('hello')

    def timeget(self):
        timeStamp = time.time()
        timeArray = time.localtime(timeStamp)
        otherStyleTime = time.strftime("%Y%m%d%H%M%S", timeArray) 
        return otherStyleTime

    def mkdir(self,path):
        path=path.strip()
        path=path.rstrip("\\")
        isExists=os.path.exists(path)
        if not isExists:
            os.makedirs(path.decode('utf-8'))
            print path + u' 创建成功'
            return path
        else:
            print path + u' 目录已存在'
            return path

    

    def CreatePptNum(self,num):
        #path = "D:\\performance2\\Result\\ppt"
        o.mkdir("D:\\performance2\\Result\\ppt")
        i = 0
        z = 0
        # 创建幻灯片 ------
        prs = Presentation()
        while i < 10:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # 定义图表数据 ---------------------
            chart_data = ChartData()
            chart_data.categories = ['East', 'West', 'Midwest']
            chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
         
            # 将图表添加到幻灯片 --------------------
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            i = i + 1

        #title_slide_layout1 = prs.slide_layouts[0]
        #slide1 = prs.slides.add_slide(title_slide_layout1)
         
        #title1 = slide1.shapes.title
        #subtitle1 = slide1.placeholders[1]
         
        # 设置标题和副标题
        #title1.text = "我是关键字"
        #salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #subtitle1.text = "王宝强-苏绍辉" + salt


        s = random.randint(0,50)
        while z < s:
            title_slide_layout1 = prs.slide_layouts[0]
            slide1 = prs.slides.add_slide(title_slide_layout1)
         
            title1 = slide1.shapes.title
            subtitle1 = slide1.placeholders[1]
            faker = Faker("zh_CN")
            text = faker.text(max_nb_chars=2048)#之前是2048
            subtitle1.text = text
            title1.text = "我是关键字"
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            subtitle1.text = "参考资料" + salt
            z = z + 1

        prs.save('D:\\performance2\\Result\\ppt\\dlp%s.pptx' % num)

        return True

    def CreatePptE(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreatePptNum(m)
            m = m + 1
        return True
    def CreatePptE1(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreatePptNum(m)
            m = m + 1
        return True

    def logtodis(self,text):
        logpath = 'D:\\performance2\\log\\rundiscovery.txt'
        f = open("D:\\performance2\\log\\rundiscovery.txt",'a+')
        f.write(text)
        f.write('---rundiscovery:\t%s\r\n' %datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
        f.close()

    def rundiscovery(self):
        o = office()
        path = "D:\\performance2\\Result\\ppt"
        isExists = os.path.exists("D:\\performance2\\Result\\ppt")
        if not isExists:
            o.mkdir("D:\\performance2\\Result\\ppt")
        else:
            shutil.rmtree("D:\\performance2\\Result\\ppt")
            time.sleep(3)
            o.mkdir("D:\\performance2\\Result\\ppt")
        
        print "START------>" + datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        o.CreatePptE1(40)#之前是40
        time.sleep(5)
        o.logtodis("PPT")
        
        print "STOP------>" + datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    def test(self):
        i = 0
        path = os.getcwd()+ '\\'
        time.sleep(3)
        starttime = time.time()
        for y in range(18):
            try:
                o = office()
                o.rundiscovery()
                time.sleep(3)
                f = o.timeget()
                print 'rundiscovery' + '  '+ str(i)+ '  ' + f

                i = i + 1
            except Exception,ex:
                print ex
            if y == 17:
                break
            time.sleep(3600)
        print 'over'
        win32api.keybd_event(win32con.VK_LWIN,0,0,0)
        win32api.keybd_event(68,0,0,0)
        win32api.keybd_event(68,0,win32con.KEYEVENTF_KEYUP,0)
        win32api.keybd_event(win32con.VK_LWIN,0,win32con.KEYEVENTF_KEYUP,0)          
        subprocess.Popen('taskmgr')
        time.sleep(4)
        im = ImageGrab.grab()
        im.save('c:\\12.png')      
        time.sleep(120)
        try:
            os.system('taskkill /IM 360tray.exe /F')
            os.system('taskkill /IM 360EntClient.exe /F')
            os.system('taskkill /IM DLPMain.exe /F')
            os.system('taskkill /IM DlpCenter.exe /F')
            os.system('taskkill /IM dlpmain32.exe /F')
            os.system('taskkill /IM DriverControler.exe /F')
            os.system('taskkill /IM DriverControler.exe /F')
        except Exception,ex:
            print ex
    
#end--------------------------------------------------------

if __name__ == "__main__":
    o = office()
    #o.test()
    o.rundiscovery()